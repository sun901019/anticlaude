from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("projects/anticlaude/interview_presentation_20260325.pptx")

ACCENT = RGBColor(220, 93, 26)
ACCENT_SOFT = RGBColor(255, 237, 220)
TEXT = RGBColor(22, 22, 22)
MUTED = RGBColor(95, 90, 83)
BG = RGBColor(250, 246, 240)
PANEL = RGBColor(255, 253, 248)
LINE = RGBColor(230, 223, 214)
SUCCESS = RGBColor(17, 117, 84)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, eyebrow, title, subtitle=None):
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.45), Inches(1.7), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_SOFT
    badge.line.color.rgb = ACCENT_SOFT
    p = badge.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = eyebrow.upper()
    r.font.name = "Aptos"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(11.5), Inches(1.45))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TEXT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.62), Inches(2.05), Inches(11.0), Inches(0.9))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.color.rgb = MUTED


def add_bullets(slide, x, y, w, h, bullets):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
        p.bullet = True
    return box


def add_panel(slide, x, y, w, h, title, body_lines):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = TEXT
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
        p.bullet = True
    return shape


def add_metric(slide, x, y, w, h, value, label, value_color=TEXT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = value_color
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = label
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER
    return shape


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = []

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(
        slide,
        "Interview Portfolio",
        "AntiClaude / AITOS / Flow Lab",
        "AI Operating System + Ecommerce Decision Workspace",
    )
    sub = slide.shapes.add_textbox(Inches(0.65), Inches(2.55), Inches(9.8), Inches(1.3))
    p = sub.text_frame.paragraphs[0]
    p.text = (
        "A full-stack system that combines workflow control, approvals, AI collaboration, "
        "content operations, and Flow Lab product economics in one operator workspace."
    )
    p.font.name = "Aptos"
    p.font.size = Pt(18)
    p.font.color.rgb = MUTED
    add_metric(slide, Inches(0.7), Inches(5.3), Inches(2.2), Inches(1.2), "351", "tests passed")
    add_metric(slide, Inches(3.0), Inches(5.3), Inches(2.2), Inches(1.2), "16", "dashboard routes")
    add_metric(slide, Inches(5.3), Inches(5.3), Inches(2.2), Inches(1.2), "101", "API handlers")
    add_metric(slide, Inches(7.6), Inches(5.3), Inches(2.2), Inches(1.2), "1262", "workflow runs")
    slides.append(slide)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Problem", "Why I built it")
    add_bullets(
        slide,
        Inches(0.8), Inches(1.8), Inches(7.4), Inches(3.7),
        [
            "AI generation had no governance or approval discipline.",
            "Ecommerce pricing and sourcing relied on ad hoc spreadsheets.",
            "Content, review, notifications, and operations were fragmented across tools.",
            "I wanted one controllable operator system instead of disconnected surfaces.",
        ],
    )
    add_panel(
        slide,
        Inches(8.6), Inches(1.9), Inches(3.8), Inches(2.25),
        "Goal",
        [
            "Turn work into workflows",
            "Make AI actions auditable",
            "Centralize operator decisions",
        ],
    )
    add_panel(
        slide,
        Inches(8.6), Inches(4.45), Inches(3.8), Inches(1.6),
        "Result",
        ["An internal operating system, not another chat app"],
    )
    slides.append(slide)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Build", "Three connected products")
    add_panel(slide, Inches(0.8), Inches(1.9), Inches(3.9), Inches(2.5), "AITOS / AI Office", [
        "AI routing",
        "approvals and review queue",
        "CEO decision packages",
        "workflow control",
    ])
    add_panel(slide, Inches(4.95), Inches(1.9), Inches(3.9), Inches(2.5), "Media Engine", [
        "topic strategy",
        "draft generation",
        "GEO / SEO related enforcement",
        "content review flow",
    ])
    add_panel(slide, Inches(9.1), Inches(1.9), Inches(3.4), Inches(2.5), "Flow Lab", [
        "product intake",
        "landed-cost pricing",
        "QQL sourcing",
        "bundle strategy",
    ])
    add_panel(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.35), "Stack", [
        "FastAPI, Next.js 14 App Router, SQLite, pytest, workflow graph runtime, AI task routing"
    ])
    slides.append(slide)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Architecture", "Backend authority + workflow runtime + operator UI")
    add_panel(slide, Inches(0.8), Inches(2.0), Inches(3.0), Inches(2.4), "API layer", [
        "chat, review, workflows, ecommerce, Flow Lab, reports, system health"
    ])
    add_panel(slide, Inches(4.1), Inches(2.0), Inches(3.0), Inches(2.4), "Runtime layer", [
        "workflow runs, tasks, approvals, checkpoints, resume / pause"
    ])
    add_panel(slide, Inches(7.4), Inches(2.0), Inches(2.9), Inches(2.4), "Domain layer", [
        "media", "flow_lab", "cleaner business logic boundaries"
    ])
    add_panel(slide, Inches(10.55), Inches(2.0), Inches(1.95), Inches(2.4), "UI", [
        "dashboard-first", "operator workflow surfaces"
    ])
    add_panel(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.1), "Architectural principle", [
        "Move from prompts and pages to governed workflows and domain logic."
    ])
    slides.append(slide)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Workflow", "AI actions are governed, not opaque")
    add_bullets(slide, Inches(0.8), Inches(1.9), Inches(6.6), Inches(3.6), [
        "workflow runs, tasks, events, artifacts",
        "approval requests and review queue",
        "CEO package generation for operator decisions",
        "pause / resume / checkpoint execution model",
    ])
    add_panel(slide, Inches(7.9), Inches(2.0), Inches(4.4), Inches(3.2), "Impact", [
        "Human-in-the-loop control",
        "Traceable AI decisions",
        "Safer high-risk actions",
        "Better observability",
    ])
    slides.append(slide)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Flow Lab", "From calculator to decision engine")
    add_panel(slide, Inches(0.8), Inches(2.0), Inches(5.45), Inches(2.8), "Bottom-Up pricing", [
        "Start from procurement and landed cost",
        "Apply Shopee fee model and role margin targets",
        "Derive recommended selling price",
    ])
    add_panel(slide, Inches(6.55), Inches(2.0), Inches(5.7), Inches(2.8), "Top-Down sourcing", [
        "Start from target market price",
        "Reverse-calculate max acceptable procurement cost",
        "Use it as a sourcing / negotiation ceiling",
    ])
    add_panel(slide, Inches(0.8), Inches(5.15), Inches(11.45), Inches(1.1), "Operator surfaces", [
        "Quick Add, product detail drawer, inbound modal, weekly performance, bundles"
    ])
    slides.append(slide)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Reliability", "I fixed trust-breaking failures, not only features")
    add_panel(slide, Inches(0.8), Inches(2.0), Inches(3.7), Inches(2.3), "Notification spam", [
        "Durable SQLite-backed dedup for repeated awaiting-human alerts"
    ])
    add_panel(slide, Inches(4.8), Inches(2.0), Inches(3.7), Inches(2.3), "Publish leak", [
        "Contained test-side social publish path so test text cannot leak to live behavior"
    ])
    add_panel(slide, Inches(8.8), Inches(2.0), Inches(3.45), Inches(2.3), "Bundle runtime repair", [
        "Aligned schema queries with real tables and restored route stability"
    ])
    add_metric(slide, Inches(1.2), Inches(5.15), Inches(4.6), Inches(1.05), "Build clean", "frontend production build passing", SUCCESS)
    add_metric(slide, Inches(6.0), Inches(5.15), Inches(5.0), Inches(1.05), "Test green", "351 passed, 1 skipped, 1 warning", SUCCESS)
    slides.append(slide)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Impact", "Why this project matters in interview terms")
    add_bullets(slide, Inches(0.8), Inches(2.0), Inches(7.0), Inches(3.9), [
        "I can bridge product thinking, backend logic, frontend UX, and reliability work.",
        "I turned practical ecommerce heuristics into software logic and operator tooling.",
        "I treated AI as a governed system, not just prompt outputs.",
        "I repeatedly repaired runtime failures and rebuilt engineering trust with tests.",
    ])
    add_panel(slide, Inches(8.25), Inches(2.1), Inches(4.0), Inches(2.8), "Good framing", [
        "Not just features",
        "Not just UI",
        "An evolving operator platform with measurable system behavior",
    ])
    slides.append(slide)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Current State", "Already usable, still improving")
    add_panel(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(3.0), "Working now", [
        "dashboard build passes",
        "backend tests pass",
        "review / approval system working",
        "Flow Lab pricing and sourcing flows working",
        "bundle suggestion route repaired",
    ])
    add_panel(slide, Inches(6.7), Inches(2.0), Inches(5.6), Inches(3.0), "Next direction", [
        "deeper backend formula convergence",
        "family / variant modeling",
        "smarter bundle recommendations",
        "further operator UX cleanup",
    ])
    slides.append(slide)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Closing", "A system I built, debugged, and kept improving")
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.2), Inches(2.2))
    p = box.text_frame.paragraphs[0]
    p.text = (
        "I built an internal AI operating system and ecommerce decision platform that is already "
        "working, measurable, and continuously refined through architecture cleanup, reliability fixes, "
        "and operator-focused UX improvements."
    )
    p.font.name = "Aptos"
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER
    add_metric(slide, Inches(3.0), Inches(5.2), Inches(2.2), Inches(1.05), "FastAPI", "backend")
    add_metric(slide, Inches(5.55), Inches(5.2), Inches(2.2), Inches(1.05), "Next.js", "frontend")
    add_metric(slide, Inches(8.1), Inches(5.2), Inches(2.2), Inches(1.05), "SQLite", "runtime data")

    prs.save(OUT)


if __name__ == "__main__":
    build()
