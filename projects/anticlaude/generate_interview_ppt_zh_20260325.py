from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("projects/anticlaude/interview_presentation_zh_20260325.pptx")

ACCENT = RGBColor(216, 92, 29)
ACCENT_SOFT = RGBColor(255, 237, 220)
TEXT = RGBColor(22, 22, 22)
MUTED = RGBColor(93, 88, 79)
BG = RGBColor(250, 246, 240)
PANEL = RGBColor(255, 253, 248)
LINE = RGBColor(230, 223, 214)
SUCCESS = RGBColor(15, 118, 110)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, eyebrow, title, subtitle=None):
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.45), Inches(1.9), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_SOFT
    badge.line.color.rgb = ACCENT_SOFT
    p = badge.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = eyebrow
    r.font.name = "Microsoft JhengHei"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(11.8), Inches(1.45))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft JhengHei UI"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TEXT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.62), Inches(2.05), Inches(11.0), Inches(0.9))
        p = sub.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Microsoft JhengHei"
        r.font.size = Pt(14)
        r.font.color.rgb = MUTED


def add_bullets(slide, x, y, w, h, bullets):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
        p.bullet = True


def add_panel(slide, x, y, w, h, title, body_lines):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft JhengHei UI"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = TEXT
    for line in body_lines:
      p = tf.add_paragraph()
      p.text = line
      p.font.name = "Microsoft JhengHei"
      p.font.size = Pt(14)
      p.font.color.rgb = MUTED
      p.bullet = True


def add_metric(slide, x, y, w, h, value, label, value_color=TEXT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = value_color
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = label
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "面試作品集", "AntiClaude / AITOS / Flow Lab", "AI 作業系統 + 電商決策工作台")
    sub = slide.shapes.add_textbox(Inches(0.65), Inches(2.55), Inches(10.2), Inches(1.3))
    p = sub.text_frame.paragraphs[0]
    p.text = "我把 AI 協作、workflow、approval、內容營運與 Flow Lab 電商決策整合成一套真正可操作的全端系統。"
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(18)
    p.font.color.rgb = MUTED
    add_metric(slide, Inches(0.7), Inches(5.3), Inches(2.2), Inches(1.2), "351", "後端測試通過")
    add_metric(slide, Inches(3.0), Inches(5.3), Inches(2.2), Inches(1.2), "16", "Dashboard 路由")
    add_metric(slide, Inches(5.3), Inches(5.3), Inches(2.2), Inches(1.2), "101", "API handlers")
    add_metric(slide, Inches(7.6), Inches(5.3), Inches(2.2), Inches(1.2), "1262", "workflow runs")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "問題", "我想解決什麼")
    add_bullets(slide, Inches(0.8), Inches(1.8), Inches(7.0), Inches(3.8), [
        "AI 可以生成很多東西，但沒有治理與審核機制。",
        "電商定價、採購、選品常依賴手算或試算表。",
        "內容、通知、審核、操作台散落在不同工具中。",
        "我希望做出一套可控、可追蹤、可審核的 operator system。",
    ])
    add_panel(slide, Inches(8.4), Inches(2.0), Inches(4.0), Inches(2.4), "專案目標", [
        "整合 AI、workflow、approval、電商決策",
        "讓系統成為真實工作台，而不是 demo 頁面",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "內容", "我做了哪些東西")
    add_panel(slide, Inches(0.8), Inches(1.9), Inches(3.9), Inches(2.5), "AITOS / AI Office", [
        "AI routing", "review queue", "approval gates", "CEO decision package"
    ])
    add_panel(slide, Inches(4.95), Inches(1.9), Inches(3.9), Inches(2.5), "Media Engine", [
        "topic strategy", "draft generation", "GEO / SEO", "content review"
    ])
    add_panel(slide, Inches(9.1), Inches(1.9), Inches(3.4), Inches(2.5), "Flow Lab", [
        "商品建檔", "落地成本", "Shopee 定價", "組合設計"
    ])
    add_panel(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.35), "技術棧", [
        "FastAPI、Next.js 14、SQLite、pytest、workflow graph runtime、AI task routing"
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "架構", "從 prompt 和頁面，收斂成 workflow 與 domain logic")
    add_panel(slide, Inches(0.8), Inches(2.0), Inches(3.0), Inches(2.4), "API 層", [
        "chat, review, workflows, ecommerce, Flow Lab, reports"
    ])
    add_panel(slide, Inches(4.1), Inches(2.0), Inches(3.0), Inches(2.4), "Runtime 層", [
        "workflow runs, tasks, approvals, checkpoints, resume / pause"
    ])
    add_panel(slide, Inches(7.4), Inches(2.0), Inches(2.9), Inches(2.4), "Domain 層", [
        "media", "flow_lab", "商業邏輯逐步分離"
    ])
    add_panel(slide, Inches(10.55), Inches(2.0), Inches(1.95), Inches(2.4), "UI", [
        "dashboard-first", "operator workflow"
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Flow Lab", "把電商從試算表，做成決策引擎")
    add_panel(slide, Inches(0.8), Inches(2.0), Inches(5.45), Inches(2.8), "Bottom-Up 定價", [
        "從採購與落地成本出發",
        "套用 Shopee 費率與角色毛利",
        "推回建議售價",
    ])
    add_panel(slide, Inches(6.55), Inches(2.0), Inches(5.7), Inches(2.8), "Top-Down 採購天花板", [
        "從市場價出發",
        "反推出最大可接受採購成本",
        "幫助議價與選品",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "可靠性", "我不只做功能，也修真正會傷害信任的問題")
    add_panel(slide, Inches(0.8), Inches(2.0), Inches(3.7), Inches(2.3), "通知 spam", [
        "SQLite-backed dedup 擋掉重複 LINE 通知"
    ])
    add_panel(slide, Inches(4.8), Inches(2.0), Inches(3.7), Inches(2.3), "測試 publish 洩漏", [
        "避免測試字串碰到真實社群發文路徑"
    ])
    add_panel(slide, Inches(8.8), Inches(2.0), Inches(3.45), Inches(2.3), "Bundle runtime", [
        "修 schema/query mismatch，讓組合建議恢復可用"
    ])
    add_metric(slide, Inches(2.0), Inches(5.2), Inches(3.8), Inches(1.05), "Build 綠燈", "前端 production build passing", SUCCESS)
    add_metric(slide, Inches(6.2), Inches(5.2), Inches(4.8), Inches(1.05), "Tests 綠燈", "351 passed, 1 skipped, 1 warning", SUCCESS)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "成果", "這個作品代表我具備跨層整合能力")
    add_bullets(slide, Inches(0.8), Inches(1.9), Inches(7.1), Inches(3.9), [
        "能把模糊需求整理成產品與流程規格",
        "能做 backend API、schema、workflow、domain logic",
        "能做 frontend operator UX 收斂",
        "能做 AI 治理與 human-in-the-loop 設計",
        "能修 runtime 問題並把測試基線救回來",
    ])
    add_panel(slide, Inches(8.35), Inches(2.2), Inches(4.0), Inches(2.6), "面試 framing", [
        "不是只做某一層",
        "而是能串起產品、實作、驗證與維護",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "下一步", "現在已可用，下一步是更成熟")
    add_panel(slide, Inches(0.8), Inches(2.0), Inches(3.7), Inches(2.3), "後端單一真相", [
        "更完整收斂商業公式到 backend"
    ])
    add_panel(slide, Inches(4.8), Inches(2.0), Inches(3.7), Inches(2.3), "Family / Variant", [
        "更完整商品家族與變體管理"
    ])
    add_panel(slide, Inches(8.8), Inches(2.0), Inches(3.45), Inches(2.3), "Bundle Intelligence", [
        "從簡單推薦升級到場景與在售商品導向"
    ])
    box = slide.shapes.add_textbox(Inches(0.9), Inches(5.15), Inches(11.2), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    p.text = "我做的是一套已經能運作、也能持續演進的 AI operator system。"
    p.font.name = "Microsoft JhengHei UI"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER

    prs.save(OUT)


if __name__ == "__main__":
    build()
